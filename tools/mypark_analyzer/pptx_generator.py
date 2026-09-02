# -*- coding: utf-8 -*-
"""McKinsey Classic Executive Theme PPTX 프레젠테이션 생성기 (PART 2 흐름 재구성 완료본)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .config import DEFAULT_SETTINGS, fmt_eok, fmt_won_full, fmt_months, DRIVING_RANGE_BENCHMARK
from .finance_engine import FinanceEngine


class PPTXGenerator:
    """16:9 와이드스크린 맥킨지 클래식 이그제큐티브 PPTX 생성기"""

    def __init__(self, filename="mypark_market_analysis.pptx"):
        self.filename = filename
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]

        # McKinsey Classic Color Palette
        self.c_mck_navy = RGBColor(0x14, 0x18, 0x1F)   # Ink #14181F
        self.c_mck_teal = RGBColor(0x1F, 0x5A, 0x44)   # Emerald accent #1F5A44
        self.c_charcoal = RGBColor(0x14, 0x18, 0x1F)   # Body text
        self.c_slate = RGBColor(0x6B, 0x6F, 0x76)      # Subtext
        self.c_line = RGBColor(0xD3, 0xD1, 0xCB)       # Divider
        self.c_box_bg = RGBColor(0xEB, 0xEA, 0xE5)     # Card bg
        self.c_tint_blue = RGBColor(0xE3, 0xEC, 0xE7)  # Emerald-tinted card bg
        self.c_white = RGBColor(255, 255, 255)
        self.c_red = RGBColor(0xB2, 0x3A, 0x2E)
        self.c_paper = RGBColor(0xF4, 0xF3, 0xF0)      # Slide background (ledger paper)

    def _add_mckinsey_header(self, slide, section_title, lead_text):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.c_paper
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.4))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.c_mck_navy
        top_bar.line.fill.background()

        tf_bar = top_bar.text_frame
        tf_bar.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_bar = tf_bar.paragraphs[0]
        p_bar.text = "   MYPARK SCREEN PARK GOLF  |  EXECUTIVE SITE SELECTION & INVESTMENT ANALYSIS"
        p_bar.font.name = 'Malgun Gothic'
        p_bar.font.size = Pt(9.5)
        p_bar.font.bold = True
        p_bar.font.color.rgb = self.c_white

        tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.48), Inches(12.133), Inches(0.42))
        tf_t = tb_title.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = section_title
        p_t.font.name = 'Malgun Gothic'
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.color.rgb = self.c_mck_navy

        tb_lead = slide.shapes.add_textbox(Inches(0.6), Inches(0.92), Inches(12.133), Inches(0.35))
        tf_l = tb_lead.text_frame
        tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
        p_l = tf_l.paragraphs[0]
        p_l.text = lead_text
        p_l.font.name = 'Malgun Gothic'
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = self.c_slate

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.30), Inches(12.133), Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = self.c_line
        line.line.fill.background()

    def _add_source_footer(self, slide, source_text="KOSIS & MYPARK Regional Tier Estimation Model"):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.05), Inches(12.133), Inches(0.012))
        line.fill.solid()
        line.fill.fore_color.rgb = self.c_line
        line.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(12.133), Inches(0.3))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"Source: {source_text}  |  CONFIDENTIAL  |  MYPARK HQ"
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(8)
        p.font.color.rgb = self.c_slate

    def _format_cell(self, cell, text, font_size=10, bold=False, color=None, bg_color=None, align=PP_ALIGN.CENTER):
        if color is None:
            color = self.c_charcoal
        if bg_color is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.alignment = align
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color

    def generate(self, data, output_pptx_path=None, charts=None):
        if output_pptx_path and isinstance(output_pptx_path, str):
            self.filename = output_pptx_path
        if charts is None:
            charts = data.get('charts', {})
        site = data['site']
        demo = data['demographics']
        comm = data['commercial']
        score = data.get('score', data.get('scores', {}))
        fin = data['financials']
        inv = fin['investment']
        scenarios = fin['monthly_scenarios']
        target_dong = site['dong']
        pop_source_tag = "MYPARK 추정 모델" if demo.get('is_estimated') else "KOSIS 실측"

        # ---------------------------------------------------------------------
        # Slide 1: 표지
        # ---------------------------------------------------------------------
        s1 = self.prs.slides.add_slide(self.blank_layout)
        bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, self.prs.slide_height)
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = self.c_mck_navy
        bg1.line.fill.background()

        tbar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.2))
        tbar.fill.solid()
        tbar.fill.fore_color.rgb = self.c_mck_teal
        tbar.line.fill.background()

        tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.8))
        tf1 = tb1.text_frame
        tf1.word_wrap = True
        
        p1 = tf1.paragraphs[0]
        p1.text = "스크린 파크골프 (마이파크) 출점 타당성 분석 보고서"
        p1.font.name = 'Malgun Gothic'
        p1.font.size = Pt(30)
        p1.font.bold = True
        p1.font.color.rgb = self.c_white

        p2 = tf1.add_paragraph()
        p2.space_before = Pt(12)
        p2.text = f"{site['rooms']}타석 {site['area_pyeong']}평 표준 모델  |  상권 분석 및 투자 타당성 평가"
        p2.font.name = 'Malgun Gothic'
        p2.font.size = Pt(17)
        p2.font.bold = True
        p2.font.color.rgb = self.c_mck_teal

        p3 = tf1.add_paragraph()
        p3.space_before = Pt(20)
        notes_txt = f"  |  특이사항: {site['special_notes']}" if site.get('special_notes') else ""
        p3.text = f"대상 주소: {site['full_address']}{notes_txt}  |  표준 모델: {site['rooms']}타석 ({site['area_pyeong']}평)"
        p3.font.name = 'Malgun Gothic'
        p3.font.size = Pt(13)
        p3.font.color.rgb = RGBColor(0xC7, 0xCB, 0xC3)

        p4 = tf1.add_paragraph()
        p4.space_before = Pt(6)
        _scope_txt = f"{demo.get('district_scope_name') or site['sigungu']} 전체 (관할 행정동 {demo.get('district_dong_count', 0)}개 전수)" if demo.get('district_wide_analysis') else f"{site['sido']} {site['sigungu']} {target_dong} 반경 3km 생활권"
        p4.text = f"상권 분석 범위: {_scope_txt}  |  분석 기준: {data.get('created_at', '2026.08')}"
        p4.font.name = 'Malgun Gothic'
        p4.font.size = Pt(13)
        p4.font.color.rgb = RGBColor(0xC7, 0xCB, 0xC3)

        tb1_bot = s1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.8))
        tf1_b = tb1_bot.text_frame
        p_b1 = tf1_b.paragraphs[0]
        p_b1.text = "마이파크(MYPARK) 가맹본부 데이터전략실"
        p_b1.font.name = 'Malgun Gothic'
        p_b1.font.size = Pt(11)
        p_b1.font.bold = True
        p_b1.font.color.rgb = self.c_white

        p_b2 = tf1_b.add_paragraph()
        p_b2.space_before = Pt(4)
        p_b2.text = "CONFIDENTIAL: 본 문서는 사업성 검토 목적 외 무단 복제 및 배포를 엄격히 금합니다."
        p_b2.font.name = 'Malgun Gothic'
        p_b2.font.size = Pt(9.5)
        p_b2.font.color.rgb = RGBColor(0x9B, 0xA7, 0x9E)

        # ---------------------------------------------------------------------
        # Slide 2: 1. 입지 적합성 종합 판정 (5-Dimension Diamond Scoring)
        # [PART 2 재구성: 재무 금액 배제, 순수 입지 평가 전진 배치]
        # ---------------------------------------------------------------------
        s3 = self.prs.slides.add_slide(self.blank_layout)
        _dw3 = demo.get('district_wide_analysis', False)
        _s3_title = "1. 구 전체 인구 및 타겟 연령 분석" if _dw3 else "1. 3km 생활권 인구 및 타겟 연령 분석"
        _s3_scope = f"{demo.get('district_scope_name') or site['sigungu']} 전체 관할" if _dw3 else "사업지 반경 3km 내"
        self._add_mckinsey_header(s3, _s3_title, f"{_s3_scope} {demo['total_pop']/10000:.1f}만 명({len(demo['dongs'])}개 행정동) 및 50대 이상 시니어 {demo['senior_50_plus']/10000:.1f}만 명({demo['senior_ratio']}%) 확보")
        
        # 좌측 행정동별 인구 테이블
        rows2 = min(len(demo['dongs']) + 2, 8)
        cols2 = 4
        left2 = Inches(0.6)
        top2 = Inches(1.45)
        width2 = Inches(5.9)
        height2 = Inches(2.72)
        table_s2 = s3.shapes.add_table(rows2, cols2, left2, top2, width2, height2).table
        table_s2.columns[0].width = Inches(1.7)
        table_s2.columns[1].width = Inches(1.4)
        table_s2.columns[2].width = Inches(1.4)
        table_s2.columns[3].width = Inches(1.4)
        
        headers2 = ["행정동명", "남성인구", "여성인구", "총 인구수"]
        for c_idx, h in enumerate(headers2):
            self._format_cell(table_s2.cell(0, c_idx), h, font_size=10, bold=True, color=self.c_white, bg_color=self.c_mck_navy)
            
        for r_idx, d in enumerate(demo['dongs'][:rows2-2]):
            r = r_idx + 1
            bg_c = self.c_tint_blue if r % 2 == 1 else self.c_white
            self._format_cell(table_s2.cell(r, 0), d['dong'], font_size=9.5, bold=True, color=self.c_mck_navy, bg_color=bg_c)
            self._format_cell(table_s2.cell(r, 1), f"{d['male']:,}", font_size=9.5, color=self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s2.cell(r, 2), f"{d['female']:,}", font_size=9.5, color=self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s2.cell(r, 3), f"{d['total']:,}", font_size=9.5, color=self.c_charcoal, bg_color=bg_c)
            
        last_r2 = rows2 - 1
        self._format_cell(table_s2.cell(last_r2, 0), "합계 (3km 생활권)", font_size=10, bold=True, color=self.c_mck_navy, bg_color=self.c_tint_blue)
        self._format_cell(table_s2.cell(last_r2, 1), f"{demo['male_pop']:,}", font_size=10, bold=True, color=self.c_mck_navy, bg_color=self.c_tint_blue)
        self._format_cell(table_s2.cell(last_r2, 2), f"{demo['female_pop']:,}", font_size=10, bold=True, color=self.c_mck_navy, bg_color=self.c_tint_blue)
        self._format_cell(table_s2.cell(last_r2, 3), f"{demo['total_pop']:,}", font_size=10, bold=True, color=self.c_mck_navy, bg_color=self.c_tint_blue)

        # 우측 시니어 세분화 매트릭스
        c3_1 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(2.72))
        c3_1.fill.solid()
        c3_1.fill.fore_color.rgb = self.c_box_bg
        c3_1.line.color.rgb = self.c_line
        tf3_1 = c3_1.text_frame
        tf3_1.word_wrap = True
        
        p = tf3_1.paragraphs[0]
        p.text = "■ 50대 이상 시니어 연령대 분포 매트릭스"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        age_matrix = [
            ("50대 (액티브 시니어)", f"{demo['pop_50s']:,}명", f"{demo['ratio_50s']}%", "부부/동호회 주말 및 평일 야간"),
            ("60대 (은퇴 시니어)", f"{demo['pop_60s']:,}명", f"{demo['ratio_60s']}%", "평일 주간 정기 리그 핵심 주력"),
            ("70대 이상 (실버 헬스케어)", f"{demo['pop_70_plus']:,}명", f"{demo['ratio_70_plus']}%", "오전 시간대 건강 증진 친목 모임"),
            ("50대 이상 총계", f"{demo['senior_50_plus']:,}명", f"{demo['senior_ratio']}%", "★ 평일 낮 10~17시 풀가동 타겟")
        ]
        for grp, cnt, rt, beh in age_matrix:
            p_m = tf3_1.add_paragraph()
            p_m.space_before = Pt(5)
            r1 = p_m.add_run()
            r1.text = f"• {grp}: "
            r1.font.bold = True
            r1.font.size = Pt(9.2)
            r1.font.color.rgb = self.c_mck_navy
            r2 = p_m.add_run()
            r2.text = f"{cnt} ({rt}), {beh}"
            r2.font.size = Pt(9.2)
            r2.font.color.rgb = self.c_charcoal

        # 하단 시사점 카드
        c3_2 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.35), Inches(12.1), Inches(2.55))
        c3_2.fill.solid()
        c3_2.fill.fore_color.rgb = self.c_box_bg
        c3_2.line.color.rgb = self.c_line
        tf3_2 = c3_2.text_frame
        tf3_2.word_wrap = True
        
        p = tf3_2.paragraphs[0]
        p.text = "■ 3km 생활권 시니어 인구 분석 시사점"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        _senior_total3 = max(1, demo['senior_50_plus'])
        _60s_share3 = demo['pop_60s'] / _senior_total3 * 100
        _70s_share3 = demo['pop_70_plus'] / _senior_total3 * 100
        insights3 = [
            f"• 압도적인 타겟 집적도: 반경 3km 내 50대 이상 인구 {demo['senior_50_plus']:,}명({demo['senior_ratio']}%) 확보로 안정적 단골 풀 형성",
            f"• 60대 주력 고객군 {_60s_share3:.0f}%: 은퇴 후 평일 낮 시간 여유가 있는 60대가 전체 시니어 중 {_60s_share3:.0f}%를 차지하여 평일 주간 가동률 극대화",
            f"• 70대 실버 헬스케어 수요 {_70s_share3:.0f}%: 관절 부담이 없는 파크골프 특성상 부부 동반 및 시니어 커뮤니티 공간으로 정착",
            "• 주간 집중 가동 구조: 평일 낮 7시간을 주력 시간대로 삼아 일일 회전수를 안정적으로 확보"
        ]
        for ins in insights3:
            p_i = tf3_2.add_paragraph()
            p_i.space_before = Pt(4)
            p_i.text = ins
            p_i.font.size = Pt(9.5)
            p_i.font.color.rgb = self.c_charcoal

        self._add_source_footer(s3, "KOSIS National Statistics Portal" + (" (※ 행정동 추정 모델 적용)" if demo.get("is_estimated") else f" ({demo.get('base_date', '2026.08')})"))

        # ---------------------------------------------------------------------
        # Slide 4: 3. 상권 소비력 및 유동 패턴 분석
        # ---------------------------------------------------------------------
        s4 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s4, "2. 상권 소비력 및 유동 패턴 분석", f"주거지역 {comm.get('residential_pop_ratio', 93.4)}% 밀집 상권 및 유사 골프업종 상위 20% 월매출 {comm['top_20_sales']//10000:,}만원 시장 타겟팅")
        
        top20_str = f"{comm['top_20_sales']//10000:,}만원"
        # 하위 20%는 산출값이 아닌 코드 기본값이라 제거하고 상권 평균으로 대체한다.
        avg_str = f"{comm['monthly_avg_sales']//10000:,}만원"
        
        c4_1 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.45), Inches(3.9), Inches(1.78))
        c4_1.fill.solid()
        c4_1.fill.fore_color.rgb = self.c_box_bg
        c4_1.line.color.rgb = self.c_line
        tf4_1 = c4_1.text_frame
        tf4_1.word_wrap = True
        p = tf4_1.paragraphs[0]
        p.text = "■ 유사 골프업종 수익구조 비교 (MYPARK 추정)"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        p_sub = tf4_1.add_paragraph()
        p_sub.space_before = Pt(4)
        p_sub.text = f"• 상위 20% 매출: {top20_str} /월 (대형 최신 매장)\n• 상권 평균 매출: {avg_str} /월 (동일 상권 평균)\n★ 마이파크 {site['rooms']}타석 플래그십은 상위 20% 시장 점유"
        p_sub.font.size = Pt(8.8)
        p_sub.font.color.rgb = self.c_charcoal
        
        c4_2 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.32), Inches(3.9), Inches(1.78))
        c4_2.fill.solid()
        c4_2.fill.fore_color.rgb = self.c_box_bg
        c4_2.line.color.rgb = self.c_line
        tf4_2 = c4_2.text_frame
        tf4_2.word_wrap = True
        p = tf4_2.paragraphs[0]
        p.text = "■ 핵심 고객층 및 이용 패턴"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        p_sub2 = tf4_2.add_paragraph()
        p_sub2.space_before = Pt(4)
        # 인접 업종(골프 연습장) 실측 지표를 함께 실어 고객층·시간대가 겹치지 않음을
        # 근거로 보여준다. 각 수치의 적용 범위(전국/사례)를 반드시 함께 적는다.
        _bm4 = DRIVING_RANGE_BENCHMARK
        p_sub2.text = (
            f"• 요일별: 주중 {100 - comm['day_distribution']['주말평균비중']*2:.1f}% / 주말 {comm['day_distribution']['주말평균비중']*2:.1f}%\n"
            f"• 시간대: 주간(10~17시) {comm['time_distribution']['주간_10_17시_비중']}% 집중 가동\n"
            f"• 인구 특성: 주거 {comm['residential_pop_ratio']}% / 직장 {comm['workplace_pop_ratio']}%\n"
            f"• 인접 업종 비교(골프 연습장): 전국 업소당 월평균 {_bm4['national_monthly_sales_manwon']:,}만원, "
            f"{_bm4['usage_scope']} 남성 {_bm4['usage_male_ratio']}%·저녁 {_bm4['usage_evening_ratio']}% 집중\n"
            f"  → 마이파크는 시니어·주간 중심이라 고객층·시간대가 겹치지 않아 보완 관계 "
            f"({_bm4['source']}, {_bm4['base_month']} 기준)")
        p_sub2.font.size = Pt(8.8)
        p_sub2.font.color.rgb = self.c_charcoal
        
        c4_3 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.18), Inches(3.9), Inches(1.78))
        c4_3.fill.solid()
        c4_3.fill.fore_color.rgb = self.c_box_bg
        c4_3.line.color.rgb = self.c_line
        tf4_3 = c4_3.text_frame
        tf4_3.word_wrap = True
        p = tf4_3.paragraphs[0]
        p.text = "■ 소비 수준 종합 평가"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        p_sub3 = tf4_3.add_paragraph()
        p_sub3.space_before = Pt(4)
        p_sub3.text = f"• 소비 등급: {comm['spending_grade']}\n• 시니어 레저/건강 소비 여력 충분\n• 동호회 정기 예약(월 단위 선결제) 중심 안정적 가동"
        p_sub3.font.size = Pt(8.8)
        p_sub3.font.color.rgb = self.c_charcoal
        
        if 'sales_trend' in charts and os.path.exists(charts['sales_trend']):
            s4.shapes.add_picture(charts['sales_trend'], Inches(4.7), Inches(1.45), width=Inches(8.0))
            
        self._add_source_footer(s4, "MYPARK Regional Tier Estimation Model")

        # ---------------------------------------------------------------------
        # Slide 5: 4. 시니어 여가 수요 지수 및 골프 특화도
        # ---------------------------------------------------------------------
        s5 = self.prs.slides.add_slide(self.blank_layout)
        _top_ind_s5 = comm['top_growth_industries'][0]
        self._add_mckinsey_header(s5, "3. 시니어 여가 수요 지수 및 골프 특화도", f"수요 지수 1위 업종 {_top_ind_s5['name']}({_top_ind_s5['growth']}) 및 전국 평균 대비 {comm['golf_industry_density']['multiple']}배 높은 골프 특화 상권")
        
        if 'growth_radar' in charts and os.path.exists(charts['growth_radar']):
            s5.shapes.add_picture(charts['growth_radar'], Inches(0.6), Inches(1.45), width=Inches(5.9))
            
        rows5 = 6
        cols5 = 4
        left5 = Inches(6.8)
        top5 = Inches(1.45)
        width5 = Inches(5.9)
        height5 = Inches(2.72)
        table_s5 = s5.shapes.add_table(rows5, cols5, left5, top5, width5, height5).table
        table_s5.columns[0].width = Inches(0.8)
        table_s5.columns[1].width = Inches(2.1)
        table_s5.columns[2].width = Inches(1.4)
        table_s5.columns[3].width = Inches(1.6)
        
        headers5 = ["순위", "업종명", "수요 지수", "업종 특성"]
        for c_idx, h in enumerate(headers5):
            self._format_cell(table_s5.cell(0, c_idx), h, font_size=9.5, bold=True, color=self.c_white, bg_color=self.c_mck_navy)
            
        for r_idx, ind in enumerate(comm['top_growth_industries'][:5]):
            r = r_idx + 1
            bg_c = self.c_tint_blue if r == 1 else (self.c_white if r % 2 == 0 else self.c_box_bg)
            self._format_cell(table_s5.cell(r, 0), str(ind['rank']), font_size=9, bold=True, color=self.c_mck_navy, bg_color=bg_c)
            self._format_cell(table_s5.cell(r, 1), ind['name'], font_size=9, bold=(r == 1), color=self.c_mck_navy if r == 1 else self.c_charcoal, bg_color=bg_c, align=PP_ALIGN.LEFT)
            self._format_cell(table_s5.cell(r, 2), ind['growth'], font_size=9, bold=True, color=self.c_red if r == 1 else self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s5.cell(r, 3), ind['status'].split('/')[0].strip(), font_size=8.5, color=self.c_slate, bg_color=bg_c)
            
        c5_bot = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(4.35), Inches(5.9), Inches(2.55))
        c5_bot.fill.solid()
        c5_bot.fill.fore_color.rgb = self.c_box_bg
        c5_bot.line.color.rgb = self.c_line
        tf5_b = c5_bot.text_frame
        tf5_b.word_wrap = True
        
        p = tf5_b.paragraphs[0]
        p.text = "■ 골프 특화 상권 시사점 및 경쟁 우위"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        _golf_ind_s5 = next((d for d in comm['top_growth_industries'] if '골프' in d['name']), comm['top_growth_industries'][0])
        insights5 = [
            f"• 레저 스포츠 소비 상위권: {comm['top_growth_industries'][0]['name']}({comm['top_growth_industries'][0]['growth']})이 1위, {_golf_ind_s5['name']}({_golf_ind_s5['growth']})이 {_golf_ind_s5['rank']}위로 시니어 여가 업종이 성장 상위 점유",
            f"• 골프 인프라 밀집도: 전국 평균 대비 {comm['golf_industry_density']['multiple']}배 높은 골프 시설 집적으로 검증된 골프 수요층 상존",
            "• 일반 골프의 파크골프 전환: 일반 골프 비용/체력 부담을 느끼는 시니어층의 스크린 파크골프 유입 가속화",
            "• 정착 단계: 시니어 여가 문화의 핵심 종목으로 자리 잡아 가는 성장 국면"
        ]
        for ins in insights5:
            p_i = tf5_b.add_paragraph()
            p_i.space_before = Pt(4)
            p_i.text = ins
            p_i.font.size = Pt(9)
            p_i.font.color.rgb = self.c_charcoal

        self._add_source_footer(s5, "MYPARK Regional Tier Estimation Model")

        # ---------------------------------------------------------------------
        # Slide 6: 5. 경쟁 환경 및 시설 공급 갭 분석
        # ---------------------------------------------------------------------
        s6 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s6, "4. 경쟁 환경 및 시설 공급 갭 분석", comm['competitor_summary'])
        
        # 카드 개수/너비는 실제로 확인된 경쟁사 수에 맞춰 동적으로 계산한다.
        # (하드코딩된 4칸 틀에 3곳만 채워 오른쪽이 비는 문제 방지)
        comps_all = comm['competitors']
        MAX_CARDS = 5
        comps = comps_all[:MAX_CARDS]
        card_gap = Inches(0.19)
        start_x = Inches(0.6)
        avail_w_in = 12.133
        n = max(1, len(comps))
        card_w = Inches((avail_w_in - (n - 1) * 0.19) / n)

        if not comps:
            # 확인된 경쟁사가 0곳(블루오션 판정)인 경우 빈 화면 대신 실제 요약 문구를 카드 자리에 표시
            empty_box = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, Inches(1.45), Inches(avail_w_in), Inches(5.45))
            empty_box.fill.solid()
            empty_box.fill.fore_color.rgb = self.c_box_bg
            empty_box.line.color.rgb = self.c_line
            tf_e = empty_box.text_frame
            tf_e.word_wrap = True
            tf_e.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_e = tf_e.paragraphs[0]
            p_e.alignment = PP_ALIGN.CENTER
            p_e.text = comm.get('competitor_summary', '반경 3km 내 확인된 경쟁 매장이 없습니다.')
            p_e.font.bold = True
            p_e.font.size = Pt(14)
            p_e.font.color.rgb = self.c_mck_navy

        for idx, c in enumerate(comps):
            cur_x = start_x + idx * (card_w + card_gap)
            
            top_box = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, cur_x, Inches(1.45), card_w, Inches(1.1))
            top_box.fill.solid()
            top_box.fill.fore_color.rgb = self.c_mck_navy
            top_box.line.fill.background()
            tf_t = top_box.text_frame
            tf_t.word_wrap = True
            tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_t = tf_t.paragraphs[0]
            p_t.alignment = PP_ALIGN.CENTER
            p_t.text = c['name']
            p_t.font.bold = True
            p_t.font.size = Pt(11)
            p_t.font.color.rgb = self.c_white
            
            mid_box = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, cur_x, Inches(2.6), card_w, Inches(0.7))
            mid_box.fill.solid()
            mid_box.fill.fore_color.rgb = self.c_tint_blue
            mid_box.line.color.rgb = self.c_line
            tf_m = mid_box.text_frame
            tf_m.word_wrap = True
            tf_m.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_m1 = tf_m.paragraphs[0]
            p_m1.alignment = PP_ALIGN.CENTER
            if c.get('rooms', 0) > 0:
                p_m1.text = f"{c.get('rooms', 0)}타석 규모"
            elif '예시 시나리오' in c.get('status', ''):
                p_m1.text = "1호점 선점 대상"
            else:
                p_m1.text = "타석 규모 확인 예정"
            p_m1.font.bold = True
            p_m1.font.size = Pt(13)
            p_m1.font.color.rgb = self.c_mck_navy
            p_m2 = tf_m.add_paragraph()
            p_m2.alignment = PP_ALIGN.CENTER
            p_m2.space_before = Pt(2)
            p_m2.text = f"[{c.get('status', '실측완료')}] {c.get('system', '스크린 시스템')}"
            p_m2.font.size = Pt(8.5)
            p_m2.font.color.rgb = self.c_slate
            
            body_box = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, cur_x, Inches(3.35), card_w, Inches(3.65))
            body_box.fill.solid()
            body_box.fill.fore_color.rgb = self.c_box_bg
            body_box.line.color.rgb = self.c_line
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            tf_body.margin_left = tf_body.margin_right = Inches(0.12)
            tf_body.margin_top = Inches(0.12)
            
            p1 = tf_body.paragraphs[0]
            p1.text = f"■ 주소: {c['address']}"
            p1.font.size = Pt(8.8)
            p1.font.color.rgb = self.c_charcoal
            
            p2 = tf_body.add_paragraph()
            p2.space_before = Pt(5)
            p2.text = f"■ 시스템: {c['system']}"
            p2.font.size = Pt(8.8)
            p2.font.color.rgb = self.c_mck_teal
            p2.font.bold = True
            
            p3 = tf_body.add_paragraph()
            p3.space_before = Pt(5)
            if c.get('rooms', 0) > 0:
                p3.text = f"■ 보유 규모: {c['rooms']}타석 운영"
            elif '예시 시나리오' in c.get('status', ''):
                p3.text = "■ 상태: 공공데이터 미등록 (신규 개업 가능성 포함해 현장 확인 권장)"
            else:
                p3.text = "■ 보유 규모: 타석수 추가 확인 예정"
            p3.font.size = Pt(8.8)
            p3.font.color.rgb = self.c_charcoal
            
            p4 = tf_body.add_paragraph()
            p4.space_before = Pt(5)
            p4.text = f"■ 특징: {c.get('features', '-')}"
            p4.font.size = Pt(8.8)
            p4.font.color.rgb = self.c_slate
            
        _comp_summary = comm.get('competitor_summary', '')
        if '예시 시나리오' in _comp_summary:
            _comp_source = "MYPARK Competitor Database Matching (Hypothetical Scenario, Live Search Unavailable)"
        elif '소상공인시장진흥공단' in _comp_summary:
            _comp_source = "SBIZ (Small Business Market Promotion Agency) Public Data"
        elif '지도 API 실시간 검색' in _comp_summary:
            _comp_source = "Live POI Search (Kakao/TMap/Naver Cross-Verified)"
        else:
            _comp_source = "MYPARK Verified National Store Database"
        self._add_source_footer(s6, _comp_source)

        # ---------------------------------------------------------------------
        # Slide 7: 6. 사업지 개요 및 현장 출점 요건 (4대 건축·인프라 체크리스트)
        # ---------------------------------------------------------------------
        s7 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s7, "5. 사업지 개요 및 현장 출점 요건", f"{site['rooms']}타석 {site['area_pyeong']}평 규모 출점을 위한 4대 건축·인프라 현장 실측 기준")
        
        base_bullets = [
            f"• 대상 주소: {site['full_address']}",
            f"• 권장 면적: 전용 {site['area_pyeong']}평 ({site['rooms']}타석 + 카페/락커룸 최적 배치)",
            f"• 층고 기준: {site['clear_height_spec']}",
            f"• 보/배관 간섭: 센서 투사 영역 및 스윙 궤적 내 장애물 사전 실측 필수",
            f"• 권장 층수: 고객 접근성 높은 지상 2~3층 권장 (쾌적한 지하 1층 가능)",
            f"• 바닥 하중: 스크린 타석 및 키오스크 하중(300kg/㎡ 이상) 적합 여부"
        ]
        if site.get('special_notes'):
            base_bullets.insert(1, f"• 고객 특이사항: {site['special_notes']}")

        cards_s7 = [
            (Inches(0.6), Inches(1.45), Inches(5.9), Inches(2.72), "■ 공간 및 유효 층고 요건", base_bullets),
            (Inches(6.8), Inches(1.45), Inches(5.9), Inches(2.72), "■ 주차 및 차량 접근성 기준", [
                f"• 주차 요건: {site['parking_spec']}",
                f"• 고객 특성: 자차 이용 시니어 비중 80% 이상으로 편리한 진출입 필수",
                f"• 진입 여건: 램프 폭 및 회전각 여유 있는 자주식 주차장 최우선",
                f"• 도로 접면: 주요 간선도로 및 대단지 아파트 진입로 인접 우수",
                f"• 보행 동선: 대중교통(버스/지하철) 도보 5~10분 생활권 완비",
                f"• 승하차 편의: 주차장에서 매장 입구까지 단차 없는 완만한 동선"
            ]),
            (Inches(0.6), Inches(4.28), Inches(5.9), Inches(2.72), "■ 건물 편의 및 승강기 설비", [
                f"• 고객 편의: {site['accessibility_spec']}",
                f"• 계단 여건: 계단 단차가 낮거나 완만한 진입 경사로 확보 필요",
                f"• 냉난방/환기: 개별 공조 및 고성능 환기 덕트 설치 공간 확인",
                f"• 소음/진동: 상하층 타 업종 간섭 방지 방음/흡음 설계 시공",
                f"• 쾌적성: 남녀 분리 청결 화장실 및 쾌적한 로비 라운지 구축",
                f"• 장애인 편의: 엘리베이터 단차 제거 및 자동문 출입구 권장"
            ]),
            (Inches(6.8), Inches(4.28), Inches(5.9), Inches(2.72), "■ 인허가 및 건축물 용도", [
                f"• 적합 용도: {site['zoning_spec']}",
                f"• 지자체 체육시설: 체육시설의 설치·이용에 관한 법률 인허가 검토",
                f"• 소방 기준: 스프링클러, 비상유도등, 비상탈출구 완비 점검",
                f"• 전기 용량: {site['rooms']}타석 시뮬레이터 동시 가동 대비 {max(25, site['rooms']*3)}kW 이상 인입",
                f"• 정화조 용량: 일 최대 150명 이상 동시 이용 기준 충족 점검",
                f"• 행정 절차: 관할 구청 건축과 및 체육진흥과 용도 사전 협의"
            ]),
        ]
        for cx, cy, cw, ch, ctitle, clines in cards_s7:
            box = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, ch)
            box.fill.solid()
            box.fill.fore_color.rgb = self.c_box_bg
            box.line.color.rgb = self.c_line
            tf_b = box.text_frame
            tf_b.word_wrap = True
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.12)
            
            p = tf_b.paragraphs[0]
            p.text = ctitle
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.c_mck_navy
            
            for l in clines[:6]:
                p_l = tf_b.add_paragraph()
                p_l.space_before = Pt(3)
                p_l.text = l
                p_l.font.size = Pt(8.8)
                p_l.font.color.rgb = self.c_charcoal

        self._add_source_footer(s7, "Building Code & Field Inspection Checklist")

        # ---------------------------------------------------------------------
        # Slide 8: [신규] 7. 표준 투자 조건 및 사업 추진 유의사항
        # [PART 2 신설: 이 보고서에서 재무 금액이 최초로 등장하는 지점 & Caveat 명시]
        # ---------------------------------------------------------------------
        s2 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s2, "6. 입지 적합성 종합 판정 (5-Dimension Diamond Scoring)", f"앞선 5개 항목 분석을 종합한 다이아몬드 스코어링 총점 {score['total_score']}점({score['grade']}등급) - {score['grade_desc']}")

        if 'radar_score' in charts and os.path.exists(charts['radar_score']):
            s2.shapes.add_picture(charts['radar_score'], Inches(0.6), Inches(1.45), width=Inches(5.6))

        tb2 = s2.shapes.add_textbox(Inches(6.4), Inches(1.45), Inches(6.3), Inches(5.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True

        # 접근성 지표: 카카오 지도로 실제 시설 개수를 센 경우(infra_is_measured)에는
        # 실측값을, 그렇지 않으면 지역등급(4단계)별 고정값을 쓴다. 어느 쪽인지는
        # 배지와 문구로 반드시 구분해 표기한다.
        _infra = comm.get('infra', {})
        _bus_count = _infra.get('버스정류장', 30)
        _subway_info = _infra.get('지하철', '')
        _has_subway = '지하철' in _subway_info or '역세권' in _subway_info or _subway_info.endswith('역')
        if _has_subway or _bus_count >= 35:
            _access_desc = f"{comm.get('spending_grade', '')} 지역등급 기준 대중교통망 우수 추정(버스정류장 약 {_bus_count}개소 수준)"
        elif _bus_count >= 20:
            _access_desc = f"{comm.get('spending_grade', '')} 지역등급 기준 표준 수준 교통망 추정(버스정류장 약 {_bus_count}개소 수준)"
        elif _bus_count >= 10:
            _access_desc = f"{comm.get('spending_grade', '')} 지역등급 기준 교통망 보통 수준 추정(버스정류장 약 {_bus_count}개소 수준, 자차 접근 동선 병행 검토 권장)"
        else:
            _access_desc = f"{comm.get('spending_grade', '')} 지역등급 기준 자차 이용 중심 상권 추정(버스정류장 약 {_bus_count}개소 수준, 주차 확보가 핵심 변수)"
        if score.get('infra_is_measured'):
            _scope_lbl8 = f"{demo.get('district_scope_name', '')} 전역" if demo.get('district_wide_analysis') else "반경 3km"
            _access_desc = f"지도 데이터 실측({_scope_lbl8} 기준): {_infra.get('지하철', '')}"
            _bits8 = [f"{k} {_infra[k]:,}개" for k in ('병원', '은행', '교육기관', '주차장')
                      if isinstance(_infra.get(k), int)]
            if _bits8:
                _access_desc += " | " + ", ".join(_bits8)
        _access_desc += " / 실제 건물 주차면·도보거리는 본사 담당자와 현장 동행 시 정밀 산정"

        # 공간 적합성: 룸/평수 미입력 시(is_auto_estimated) 표준값 대신 중립 점수를
        # 적용한 것이므로 산출근거 문구도 그 사실을 그대로 반영한다.
        _pyeong_per_room = site['area_pyeong'] / max(1, site['rooms'])
        if not score.get('space_is_verified', True):
            _space_desc = "룸/평수 미입력으로 표준값 대신 중립 점수 적용 (현장 실측 시 정밀 산정)"
        elif _pyeong_per_room >= 12.0:
            _space_desc = f"{site['area_pyeong']}평 {site['rooms']}타석 (타석당 {_pyeong_per_room:.1f}평), 여유로운 플래그십 규모"
        elif _pyeong_per_room >= 10.0:
            _space_desc = f"{site['area_pyeong']}평 {site['rooms']}타석 (타석당 {_pyeong_per_room:.1f}평), 표준 배치 규모"
        elif _pyeong_per_room >= 8.0:
            _space_desc = f"{site['area_pyeong']}평 {site['rooms']}타석 (타석당 {_pyeong_per_room:.1f}평), 효율 배치 설계 권장"
        else:
            _space_desc = f"{site['area_pyeong']}평 {site['rooms']}타석 (타석당 {_pyeong_per_room:.1f}평), 타석 수·면적 조합을 전문가와 함께 재설계 권장"
        if score.get('space_is_verified', True):
            _space_desc += " / 권장 유효 층고 2.8m 이상 여부는 '인테리어 실측' 필수"

        indicators = [
            ("1) 시니어 인구 밀집도", score['scores']['senior_population'], 25,
             f"{pop_source_tag}: 반경 3km 내 50대 이상 {demo['senior_50_plus']:,}명({demo['senior_ratio']}%) / 본 매장 운영에 필요한 단골 약 {score.get('senior_customers_needed', 1200):,}명은 배후 시니어의 {score.get('senior_penetration', 0)}% 수준",
             not demo.get('is_estimated', False)),
            ("2) 접근성 및 주차 인프라", score['scores']['accessibility_parking'], 25, _access_desc, False),
            ("3) 공간 적합성 및 층고", score['scores']['space_efficiency'], 15, _space_desc, score.get('space_is_verified', True)),
            ("4) 경쟁 매장 여유도", score['scores']['supply_gap'], 15, f"{comm.get('competitor_summary', '반경 3km 내 대형 플래그십 매장 공급 부족')}", score.get('gap_is_verified', False)),
            ("5) 지역 소비력 및 여가지출", score['scores']['commercial_spending'], 20,
             f"MYPARK 지역등급(4단계 분류) 추정치: 시니어 여가 수요 지수 {comm.get('growth_rate', 145.2):.0f}점 및 스크린골프 상위 20% 월 {comm['top_20_sales']//10000:,}만원 상권 (동일 등급 지역은 동일 수치 적용, 개별 카드매출 실측 아님)", False),
        ]
        for idx, (iname, iscore, imax, idesc, iverified) in enumerate(indicators):
            badge_color = self.c_mck_teal if iverified else RGBColor(0x9A, 0xA5, 0xB1)
            badge_text = "실측·확인" if iverified else "추정·정밀분석 권장"
            p = tf2.add_paragraph() if idx > 0 else tf2.paragraphs[0]
            p.space_before = Pt(6)
            r1 = p.add_run()
            r1.text = f"■ {iname}: "
            r1.font.bold = True
            r1.font.size = Pt(11)
            r1.font.color.rgb = self.c_mck_navy
            r2 = p.add_run()
            r2.text = f"{iscore}점 / {imax}점 만점  ["
            r2.font.bold = True
            r2.font.size = Pt(11)
            r2.font.color.rgb = self.c_mck_teal
            r3 = p.add_run()
            r3.text = badge_text
            r3.font.bold = True
            r3.font.size = Pt(9)
            r3.font.color.rgb = badge_color
            r4 = p.add_run()
            r4.text = "]"
            r4.font.bold = True
            r4.font.size = Pt(11)
            r4.font.color.rgb = self.c_mck_teal
            p_desc = tf2.add_paragraph()
            p_desc.text = f"   ↳ 산출 근거: {idesc}"
            p_desc.font.size = Pt(9.2)
            p_desc.font.color.rgb = self.c_slate

        box2_res = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(5.5), Inches(6.3), Inches(1.4))
        box2_res.fill.solid()
        box2_res.fill.fore_color.rgb = self.c_box_bg
        box2_res.line.color.rgb = self.c_line
        tf_b2 = box2_res.text_frame
        tf_b2.word_wrap = True
        
        p_b2_t = tf_b2.paragraphs[0]
        p_b2_t.text = f"★ 종합 입지 판정: 총점 {score['total_score']}점 / {score['grade']}등급 ({score['grade_desc']})"
        p_b2_t.font.bold = True
        p_b2_t.font.size = Pt(11.5)
        p_b2_t.font.color.rgb = self.c_mck_navy
        
        grade_summary_lines_pptx = {
            'S': "본 사업지는 50~70대 풍부한 시니어 수요와 우수한 접근성을 갖추어 출점에 최적화된 입지입니다.",
            'A': "본 사업지는 시니어 배후 수요와 접근성 등 핵심 조건을 대체로 충족하는 안정적인 입지입니다.",
            'B': "본 사업지는 일부 지표가 표준 기준에 못 미쳐, 아래 세부 근거를 현장 실측과 함께 신중히 검토해야 합니다.",
            'C': "본 사업지는 5대 지표 중 다수가 표준 기준에 미달하여, 출점 전 현장 재확인이 반드시 필요합니다.",
        }
        p_b2_d = tf_b2.add_paragraph()
        p_b2_d.space_before = Pt(4)
        p_b2_d.text = f"• {grade_summary_lines_pptx.get(score['grade'], grade_summary_lines_pptx['B'])}\n• 위 5개 지표의 상세 근거는 앞선 1~5번 슬라이드(인구·상권·업종·경쟁·사업지 개요)를 참조하십시오."
        p_b2_d.font.size = Pt(9.2)
        p_b2_d.font.color.rgb = self.c_charcoal

        self._add_source_footer(s2, f"MYPARK 5-Dimension Diamond Scoring Methodology ({score['total_score']}점 {score['grade']}등급)")

        # ---------------------------------------------------------------------
        # Slide 3: 2. 배후 인구 및 타겟 연령 분석
        # ---------------------------------------------------------------------
        s8 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s8, "7. 표준 투자 조건 및 사업 추진 유의사항", f"{site['rooms']}타석 {site['area_pyeong']}평 표준 모델 기준 및 투자 결정 전 필수 점검사항")
        
        # 블록 A (좌측): 표준 투자 조건
        box8_a = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.45), Inches(5.9), Inches(5.5))
        box8_a.fill.solid()
        box8_a.fill.fore_color.rgb = self.c_box_bg
        box8_a.line.color.rgb = self.c_line
        
        top8_a = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.45), Inches(5.9), Inches(0.5))
        top8_a.fill.solid()
        top8_a.fill.fore_color.rgb = self.c_mck_navy
        top8_a.line.fill.background()
        tf8_at = top8_a.text_frame
        p8_at = tf8_at.paragraphs[0]
        p8_at.text = f" ■ {site['rooms']}타석 {site['area_pyeong']}평 표준 모델 투자 조건 (SSOT)"
        p8_at.font.bold = True
        p8_at.font.size = Pt(11)
        p8_at.font.color.rgb = self.c_white
        
        tb8_ab = s8.shapes.add_textbox(Inches(0.75), Inches(2.05), Inches(5.6), Inches(4.7))
        tf8_ab = tb8_ab.text_frame
        tf8_ab.word_wrap = True
        
        p = tf8_ab.paragraphs[0]
        p.text = "● 초기 투자금 상세 내역"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = self.c_mck_navy
        
        _rent_tag8 = "입력하신 실측" if not site.get('rent_is_estimated') else "지역 시세 추정"
        items8_a = [
            f"• 시뮬레이터 장비: {site['rooms']}대 × 대당 {DEFAULT_SETTINGS['simulator_unit_price']//10000:,}만원 = {fmt_won_full(inv['simulator_cost'])}",
            f"• 인테리어 공사비: {site['area_pyeong']}평 × 평당 {DEFAULT_SETTINGS['interior_cost_per_pyeong']//10000:,}만원 = {fmt_won_full(inv['interior_cost'])}",
            f"• 부대설비 (간판/가구/초도용품): {fmt_won_full(inv['other_facilities'])}",
            f"  - 간판({DEFAULT_SETTINGS['signage_cost']//10000:,}만) / 가구({DEFAULT_SETTINGS['furniture_cost']//10000:,}만) / 초도용품({DEFAULT_SETTINGS['supplies_cost']//10000:,}만)",
            f"• 냉난방 설비 (선택): 기존 상가 설비 승계 시 추가 비용 없음, 신규 설치 시 약 {inv['hvac_cost_optional']//10000:,}만원 별도",
            f"★ 총 초기 투자금: {fmt_won_full(inv['total_capex'])} ({fmt_eok(inv['total_capex'])})",
            "",
            "● 표준 운영 방식 및 인건비 모델",
            f"• 표준 모델 (점주 {site['staff_count']}인 상주): 인건비 월 {site['staff_count']*DEFAULT_SETTINGS['labor_cost_manager']//10000:,}만원 — 부부 공동운영 등 실제 운영 형태에 따라 조정 가능한 대표값",
            f"• 비교 모델 (직원 3인 채용): 인건비 월 {3*DEFAULT_SETTINGS['labor_cost_manager']//10000:,}만원 (회수기간 {fin['owner_operated']['staff3_payback_months']:.1f}개월)",
            f"• 게임비 요금: 1인 18홀 7,000원 (4인 1팀 28,000원)",
            f"• 3대 매출원: 게임비 회전 + 용품 판매(월 {fin['monthly_scenarios']['moderate']['goods_revenue']//10000:,}만) + 식음료(월 {fin['monthly_scenarios']['moderate']['beverage_revenue']//10000:,}만)",
            "• 본 분석 기준: 건물주(자가 소유) — 월 임대료가 발생하지 않는 조건으로 산출",
        ]
        _tp8 = data.get('tenant_payback_months')
        _tenant_txt8 = f"약 {_tp8:.1f}개월" if _tp8 else "별도 산출"
        items8_a.append(f"• 임차 운영 시 참고: {_rent_tag8} 임대료 {site['monthly_rent']//10000:,}만원/월 반영 시 회수기간 {_tenant_txt8}")
        for it in items8_a:
            p_it = tf8_ab.add_paragraph()
            p_it.space_before = Pt(3)
            p_it.text = it
            if "★ 총 초기 투자금" in it:
                p_it.font.bold = True
                p_it.font.size = Pt(11.5)
                p_it.font.color.rgb = self.c_red
            elif "●" in it:
                p_it.font.bold = True
                p_it.font.size = Pt(11)
                p_it.font.color.rgb = self.c_mck_navy
            else:
                p_it.font.size = Pt(9.2)
                p_it.font.color.rgb = self.c_charcoal

        # 블록 B (우측): 투자 결정 전 유의사항 (Caveat)
        box8_b = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(5.5))
        box8_b.fill.solid()
        box8_b.fill.fore_color.rgb = self.c_box_bg
        box8_b.line.color.rgb = self.c_line
        
        top8_b = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(0.5))
        top8_b.fill.solid()
        top8_b.fill.fore_color.rgb = self.c_mck_navy
        top8_b.line.fill.background()
        tf8_bt = top8_b.text_frame
        p8_bt = tf8_bt.paragraphs[0]
        p8_bt.text = " 참고사항 (사업 검토 시 확인해 주십시오)"
        p8_bt.font.bold = True
        p8_bt.font.size = Pt(11)
        p8_bt.font.color.rgb = self.c_white
        
        tb8_bb = s8.shapes.add_textbox(Inches(6.95), Inches(2.05), Inches(5.6), Inches(4.7))
        tf8_bb = tb8_bb.text_frame
        tf8_bb.word_wrap = True
        
        _comp_summary8 = comm.get('competitor_summary', '')
        if '예시 시나리오' in _comp_summary8:
            _methodology_txt8 = "본 보고서는 KOSIS 인구통계, MYPARK 지역등급 추정 모델, 표준 재무 모델에 기반한 추정 분석 자료입니다."
        else:
            _methodology_txt8 = "본 보고서는 인구통계, 실측 경쟁사 데이터, MYPARK 지역등급 추정 모델, 표준 재무 모델을 결합한 종합분석자료입니다."

        items8_b = [
            ("● 현장 실측 및 인허가 유의사항", [
                "• 위 수치는 표준 모델 기준 추정치이며, 실제 임대료·공사비는 현장 견적에 따라 달라질 수 있습니다.",
                "• 건물 내 보/배관 간섭 및 유효 층고(2.8m 이상), 전력 용량(30kW 이상)을 확인하십시오."
            ]),
            ("● 인건비 및 운영 방식 유의사항", [
                "• 매니저/직원을 채용해 위탁 운영할 경우 인건비 증가(월 500~750만)로 회수기간이 늘어납니다.",
                "• 인테리어 및 시뮬레이터 단가는 본 계약 시점의 공식 견적을 확인하십시오."
            ]),
            ("● 분석 방법론 및 활용 안내", [
                f"• {_methodology_txt8}",
                "• 마이파크 사업부서 전문가와의 상담을 권장하며, 특정 수익률을 보장하지 않습니다."
            ])
        ]
        for idx, (btitle, blines) in enumerate(items8_b):
            p_t = tf8_bb.add_paragraph() if idx > 0 else tf8_bb.paragraphs[0]
            p_t.space_before = Pt(6)
            p_t.text = btitle
            p_t.font.bold = True
            p_t.font.size = Pt(10.5)
            p_t.font.color.rgb = self.c_mck_navy
            for bl in blines:
                p_l = tf8_bb.add_paragraph()
                p_l.space_before = Pt(2)
                p_l.text = bl
                p_l.font.size = Pt(8.8)
                p_l.font.color.rgb = self.c_charcoal

        self._add_source_footer(s8, "MYPARK Standard Investment Criteria & Reference Notes")

        # ---------------------------------------------------------------------
        # Slide 9: 8. 사업 타당성 분석 - 매출 추정 (3대 시나리오)
        # ---------------------------------------------------------------------
        s9 = self.prs.slides.add_slide(self.blank_layout)
        m_scen = fin['monthly_scenarios']
        self._add_mckinsey_header(s9, "8. 사업 타당성 분석 - 매출 추정 (3대 시나리오)", f"보수적(일 {m_scen['conservative']['daily_turns_per_room']}회전) {m_scen['conservative']['total_revenue']//10000:,}만원 ~ 보편적(일 {m_scen['moderate']['daily_turns_per_room']}회전) {m_scen['moderate']['total_revenue']//10000:,}만원 ~ 긍정적(일 {m_scen['optimistic']['daily_turns_per_room']}회전) {m_scen['optimistic']['total_revenue']//10000:,}만원")
        
        # 1. 상단 안내 박스
        top_s9 = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.45), Inches(12.133), Inches(0.65))
        top_s9.fill.solid()
        top_s9.fill.fore_color.rgb = self.c_tint_blue
        top_s9.line.color.rgb = self.c_line
        tf_s9 = top_s9.text_frame
        tf_s9.word_wrap = True
        p_s9 = tf_s9.paragraphs[0]
        p_s9.text = f"■ {site['rooms']}타석 {site['area_pyeong']}평 기준 3대 시나리오별 월간/연간 매출 추정 (1인 18홀 {DEFAULT_SETTINGS['game_price_18hole']:,}원 / 4인 1팀 {DEFAULT_SETTINGS['game_price_18hole']*4:,}원)"
        p_s9.font.size = Pt(11)
        p_s9.font.bold = True
        p_s9.font.color.rgb = self.c_mck_navy
        
        # 2. 매출 시나리오 정밀 비교 테이블
        rows9 = 4
        cols9 = 6
        table_s9 = s9.shapes.add_table(rows9, cols9, Inches(0.6), Inches(2.25), Inches(12.133), Inches(2.8)).table
        table_s9.columns[0].width = Inches(2.2)
        table_s9.columns[1].width = Inches(1.8)
        table_s9.columns[2].width = Inches(2.0)
        table_s9.columns[3].width = Inches(2.0)
        table_s9.columns[4].width = Inches(2.1)
        table_s9.columns[5].width = Inches(2.033)
        
        headers9 = ["구분 / 시나리오", "일 회전수 (타석당)", "월 게임비 매출", "월 부가 매출(용품+식음)", "월 총매출액", "일 이용 인원"]
        for c_idx, h in enumerate(headers9):
            self._format_cell(table_s9.cell(0, c_idx), h, font_size=10, bold=True, color=self.c_white, bg_color=self.c_mck_navy)
            
        sc_keys = [('conservative', '보수적 시나리오 (비수기/초기)'), ('moderate', '보편적 시나리오 (정기예약 정착)'), ('optimistic', '긍정적 시나리오 (주간/주말 풀가동)')]
        for r_idx, (k, label) in enumerate(sc_keys):
            r = r_idx + 1
            sc = m_scen[k]
            is_mod = (k == 'moderate')
            bg_c = self.c_tint_blue if is_mod else (self.c_white if r % 2 == 0 else self.c_box_bg)
            
            self._format_cell(table_s9.cell(r, 0), label, font_size=10, bold=True, color=self.c_mck_navy if is_mod else self.c_charcoal, bg_color=bg_c, align=PP_ALIGN.LEFT)
            self._format_cell(table_s9.cell(r, 1), f"타석당 1일 {sc['daily_turns_per_room']}회전", font_size=9.5, color=self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s9.cell(r, 2), f"{sc['room_revenue']//10000:,}만원", font_size=9.5, color=self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s9.cell(r, 3), f"{(sc['goods_revenue'] + sc['fnb_revenue'])//10000:,}만원", font_size=9.5, color=self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s9.cell(r, 4), f"{sc['total_revenue']//10000:,}만원 (연 {sc['annual_revenue']//100000000:.1f}억)", font_size=10, bold=True, color=self.c_red if is_mod else self.c_mck_navy, bg_color=bg_c)
            self._format_cell(table_s9.cell(r, 5), f"1일 {sc['daily_users']}명 (월 {sc['monthly_users']:,}명)", font_size=9.5, color=self.c_slate, bg_color=bg_c)
            
        # 3. 하단 시나리오별 시사점 콜아웃 카드
        c_sc = m_scen['conservative']
        m_sc = m_scen['moderate']
        o_sc = m_scen['optimistic']
        
        callouts = [
            (Inches(0.6), f"■ 보수적 시나리오 (월 {c_sc['total_revenue']//10000:,}만원)",
             f"• 타석당 1일 {c_sc['daily_turns_per_room']}회전(1일 {c_sc['daily_users']}명) 가동 기준\n"
             f"• 월 순영업이익 약 {c_sc['operating_profit']//10000:,}만원(이익률 {c_sc['profit_margin']}%) 확보\n"
             f"• 손익분기 월매출(약 {fin['investment']['bep_monthly_sales']//10000:,}만원) 초과 안정 수익"),
            (Inches(4.68), f"■ 보편적 시나리오 (월 {m_sc['total_revenue']//10000:,}만원)",
             f"• 평일 주간 10~17시 동호회 정기 예약 중심 일 {m_sc['daily_turns_per_room']}회전 가동\n"
             f"• 월 순영업이익 {m_sc['operating_profit']//10000:,}만원(이익률 {m_sc['profit_margin']}%) 달성\n"
             f"• 단 {fin['investment']['payback_months_moderate']:.1f}개월(약 {fmt_months(fin['investment']['payback_months_moderate'])}) 만에 {fmt_eok(inv['total_capex'])} 전액 회수"),
            (Inches(8.76), f"■ 긍정적 시나리오 (월 {o_sc['total_revenue']//10000:,}만원)",
             f"• 주말 단체 예약 및 주간 풀가동 일 {o_sc['daily_turns_per_room']}회전({o_sc['daily_users']}명) 달성\n"
             f"• 월 순영업이익 {o_sc['operating_profit']//10000:,}만원(이익률 {o_sc['profit_margin']}%) 극대화\n"
             f"• 연간 순영업익 약 {o_sc['operating_profit']*12//100000000:.1f}억원 창출 플래그십 모델")
        ]
        for bx, btitle, bdesc in callouts:
            c_box = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, Inches(5.3), Inches(3.96), Inches(1.7))
            c_box.fill.solid()
            c_box.fill.fore_color.rgb = self.c_box_bg
            c_box.line.color.rgb = self.c_line
            tf_cb = c_box.text_frame
            tf_cb.word_wrap = True
            tf_cb.margin_left = tf_cb.margin_right = tf_cb.margin_top = tf_cb.margin_bottom = Inches(0.12)
            
            p_ct = tf_cb.paragraphs[0]
            p_ct.text = btitle
            p_ct.font.size = Pt(10)
            p_ct.font.bold = True
            p_ct.font.color.rgb = self.c_mck_navy
            
            p_cd = tf_cb.add_paragraph()
            p_cd.space_before = Pt(3)
            p_cd.text = bdesc
            p_cd.font.size = Pt(8.5)
            p_cd.font.color.rgb = self.c_charcoal

        self._add_source_footer(s9, "MYPARK Standard Financial Simulation Model")

        # ---------------------------------------------------------------------
        # Slide 10: 9. 사업 타당성 분석 - 비용 구조 및 순영업이익
        # ---------------------------------------------------------------------
        s10 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s10, "9. 사업 타당성 분석 - 비용 구조 및 순영업이익", f"건물주(자가 소유) 기준 월 고정비 {fin['owner_operated']['fixed_cost']//10000:,}만원(인건비 250만+운영비, 임대료 미포함) 및 보편 월 순영업이익 {m_scen['moderate']['operating_profit']//10000:,}만원")
        
        if 'waterfall_cost' in charts and os.path.exists(charts['waterfall_cost']):
            s10.shapes.add_picture(charts['waterfall_cost'], Inches(0.6), Inches(1.45), width=Inches(5.9))
            
        c10_1 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(2.65))
        c10_1.fill.solid()
        c10_1.fill.fore_color.rgb = self.c_box_bg
        c10_1.line.color.rgb = self.c_line
        tf10_1 = c10_1.text_frame
        tf10_1.word_wrap = True
        
        p = tf10_1.paragraphs[0]
        p.text = "■ 월간 비용 구조 상세 (보편 시나리오 기준)"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        _rent_basis10 = "입력하신 임대료" if not site.get('rent_is_estimated') else "지역 시세 추정 임대료"
        cost_details = [
            f"• 월 임대료: 0원 (건물주 자가 소유 기준)  ※ 임차 시 {_rent_basis10} {site['monthly_rent']//10000:,}만원 별도",
            f"• 인건비 (점주 직접운영): {m_scen['moderate']['labor_cost']//10000:,}만원 ({site['staff_count']}인 상주 운영 체제)",
            f"• 매장 운영비/소모품: {DEFAULT_SETTINGS['store_ops_monthly']//10000:,}만원  |  통신/POS: {DEFAULT_SETTINGS['pos_telecom_monthly']//10000:,}만원  |  마케팅비: {DEFAULT_SETTINGS['marketing_monthly']//10000:,}만원",
            f"• 변동비: 원가 {(m_scen['moderate']['cost_goods']+m_scen['moderate']['cost_beverage'])//10000:,}만원 + 카드수수료({DEFAULT_SETTINGS['card_fee_rate']*100:.1f}%) {m_scen['moderate']['card_fee']//10000:,}만원",
            f"★ 월 총지출 합계: {m_scen['moderate']['total_cost']//10000:,}만원",
            f"★ 월 순영업이익: {m_scen['moderate']['operating_profit']//10000:,}만원 (이익률 {m_scen['moderate']['profit_margin']}%)"
        ]
        for cd in cost_details:
            p_c = tf10_1.add_paragraph()
            p_c.space_before = Pt(3)
            p_c.text = cd
            if "★" in cd:
                p_c.font.bold = True
                p_c.font.size = Pt(10)
                p_c.font.color.rgb = self.c_mck_teal if "순영업이익" in cd else self.c_mck_navy
            else:
                p_c.font.size = Pt(9)
                p_c.font.color.rgb = self.c_charcoal
                
        c10_2 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(4.25), Inches(5.9), Inches(2.65))
        c10_2.fill.solid()
        c10_2.fill.fore_color.rgb = self.c_box_bg
        c10_2.line.color.rgb = self.c_line
        tf10_2 = c10_2.text_frame
        tf10_2.word_wrap = True
        
        p = tf10_2.paragraphs[0]
        p.text = "■ 운영 모델별 순영업이익 비교"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        models_comp = [
            f"• 점주 직접 운영 모델 (표준): 월 순영업이익 {m_scen['moderate']['operating_profit']//10000:,}만원 (연간 {m_scen['moderate']['operating_profit']*12//10000:,}만원 / 이익률 {m_scen['moderate']['profit_margin']}%)",
            f"• 직원 채용 모델 (매니저 1인 + 알바 2인): 월 순영업이익 {fin['owner_operated']['staff3_operating_profit']//10000:,}만원 (연간 {fin['owner_operated']['staff3_operating_profit']*12//10000:,}만원)",
            "• 낮은 변동비 구조: 원재료비 비중이 매우 낮아 매출이 늘수록 순이익이 빠르게 커지는 고마진 구조",
            "• 고정비 방어력: 월 고정비가 낮아 비수기나 상권 초기 단계에서도 안정적 순영업이익 기반 유지"
        ]
        for mc in models_comp:
            p_m = tf10_2.add_paragraph()
            p_m.space_before = Pt(3)
            p_m.text = mc
            p_m.font.size = Pt(9)
            p_m.font.color.rgb = self.c_charcoal

        self._add_source_footer(s10, "MYPARK Cost Structure & Operating Profit Analysis")

        # ---------------------------------------------------------------------
        # Slide 11: 10. 손익분기점(BEP) 및 투자금 회수기간
        # ---------------------------------------------------------------------
        s11 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s11, "10. 손익분기점(BEP) 및 투자금 회수기간", f"손익분기 매출 월 {inv['bep_monthly_sales']//10000:,}만원 (타석당 일 {inv['bep_turns_per_room']}회전) 및 투자금 {fmt_eok(inv['total_capex'])} 회수기간 약 {inv['payback_months_moderate']:.1f}개월")
        
        if 'bep_chart' in charts and os.path.exists(charts['bep_chart']):
            s11.shapes.add_picture(charts['bep_chart'], Inches(0.6), Inches(1.45), width=Inches(5.9))
            
        c11_1 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.9), Inches(2.65))
        c11_1.fill.solid()
        c11_1.fill.fore_color.rgb = self.c_box_bg
        c11_1.line.color.rgb = self.c_line
        tf11_1 = c11_1.text_frame
        tf11_1.word_wrap = True
        
        p = tf11_1.paragraphs[0]
        p.text = f"■ 투자금 {fmt_eok(inv['total_capex'])} 회수 시뮬레이션"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        if inv.get('conservative_viable', True):
            _bep_con_line = f"• 보수적 시나리오: 월 순익 {m_scen['conservative']['operating_profit']//10000:,}만원 -> 회수기간 약 {inv['payback_months_conservative']:.1f}개월"
        else:
            _bep_con_line = "• 보수적 시나리오(3회전)는 추가 매출 확보 전략이 필요한 구간이며, 최소 4회전 이상 가동 시 안정적 순익 구조로 전환됩니다"
        bep_sim = [
            _bep_con_line,
            f"• 보편적 시나리오: 월 순익 {m_scen['moderate']['operating_profit']//10000:,}만원 -> 회수기간 약 {inv['payback_months_moderate']:.1f}개월 (약 {fmt_months(inv['payback_months_moderate'])})",
            f"• 긍정적 시나리오: 월 순익 {m_scen['optimistic']['operating_profit']//10000:,}만원 -> 회수기간 약 {inv['payback_months_optimistic']:.1f}개월",
            f"★ BEP 달성 요건: 타석당 하루 {inv['bep_turns_per_room']}회전 (1일 {inv['bep_daily_users']}명 이용)",
            f"★ 일 평균 {inv['bep_daily_users']}명 방문 시 월 고정비 전액 커버 (보편 시나리오의 {inv['bep_daily_users']/max(1,m_scen['moderate']['daily_users'])*100:.0f}% 수준)"
        ]
        for bs in bep_sim:
            p_b = tf11_1.add_paragraph()
            p_b.space_before = Pt(3)
            p_b.text = bs
            if "★" in bs:
                p_b.font.bold = True
                p_b.font.size = Pt(9.5)
                p_b.font.color.rgb = self.c_red
            else:
                p_b.font.size = Pt(9)
                p_b.font.color.rgb = self.c_charcoal
                
        c11_2 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(4.25), Inches(5.9), Inches(2.65))
        c11_2.fill.solid()
        c11_2.fill.fore_color.rgb = self.c_box_bg
        c11_2.line.color.rgb = self.c_line
        tf11_2 = c11_2.text_frame
        tf11_2.word_wrap = True
        
        p = tf11_2.paragraphs[0]
        p.text = "■ 투자 안정성 및 리스크 평가"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.c_mck_navy
        
        risks = [
            f"• 안정적 손익 구조: 타석당 하루 {inv['bep_turns_per_room']}회전(1일 {inv['bep_daily_users']}명) 이상 가동 시 고정비 전액 커버",
            f"• 빠른 자본 회수: 보편 가동 기준 약 {inv['payback_months_moderate']:.1f}개월({fmt_months(inv['payback_months_moderate'])}) 만에 초기 투자금 {fmt_eok(inv['total_capex'])} 전액 회수",
            "• 자산 가치 보존: 시뮬레이터 장비 및 쾌적한 인테리어 시설은 향후 지속적인 현금 흐름을 창출하는 핵심 실물 자산",
            "• 안정적 단골 락인: 지역 시니어 동호회 정기 예약 시스템 구축으로 경기 변동에 영향을 받지 않는 방어적 사업 모델",
            "• 본사 오픈 지원: 마이파크 본사가 오픈 초기 홍보 및 지역 커뮤니티 형성을 함께 지원해 조기 가동률 확보를 돕습니다"
        ]
        for rk in risks:
            p_r = tf11_2.add_paragraph()
            p_r.space_before = Pt(3)
            p_r.text = rk
            p_r.font.size = Pt(9)
            p_r.font.color.rgb = self.c_charcoal

        self._add_source_footer(s11, "MYPARK BEP & Capital Payback Period Analysis")

        # ---------------------------------------------------------------------
        # Slide 12: 11. 5개년 중장기 손익 전망 및 종합 제언
        # ---------------------------------------------------------------------
        s12 = self.prs.slides.add_slide(self.blank_layout)
        self._add_mckinsey_header(s12, "11. 5개년 중장기 손익 전망 및 종합 제언", f"반경 3km 내 {demo['senior_50_plus']/10000:.1f}만 시니어 소비 수요와 주간 풀가동으로 {inv['payback_months_moderate']:.1f}개월 내 투자금 전액 회수 가능")
        
        # 1. 상단 5개년 손익 예측 테이블
        rows12 = 5
        cols12 = 7
        table_s12 = s12.shapes.add_table(rows12, cols12, Inches(0.6), Inches(1.45), Inches(12.133), Inches(1.15)).table
        table_s12.columns[0].width = Inches(1.8)
        for i in range(1, 7):
            table_s12.columns[i].width = Inches(1.72)
            
        headers12 = ["연차 구분", "1차년도", "2차년도", "3차년도", "4차년도", "5차년도", "5개년 누적"]
        for c_idx, h in enumerate(headers12):
            self._format_cell(table_s12.cell(0, c_idx), h, font_size=9.5, bold=True, color=self.c_white, bg_color=self.c_mck_navy)
            
        fy = fin['five_year']
        rows_5y = [
            ("연간 총매출", [f"{y['revenue']/100000000:.2f}억" for y in fy['years']], f"{fy['total_5yr_revenue']//100000000:.1f}억원"),
            ("연간 총비용", [f"{y['cost']/100000000:.2f}억" for y in fy['years']], f"{fy['total_5yr_cost']//100000000:.1f}억원"),
            ("연간 순영업익", [f"{y['profit']/100000000:.2f}억" for y in fy['years']], f"{fy['total_5yr_profit']//100000000:.1f}억원"),
            ("투자금 누적회수", [f"{fmt_eok(inv['total_capex'])} 회수" if fy['years'][i]['cumulative_profit'] >= inv['total_capex'] else f"{fy['years'][i]['cumulative_profit']/100000000:.2f}억" for i in range(5)], f"회수율 {fy['total_5yr_profit']/inv['total_capex']*100:.0f}%")
        ]
        for r_idx, (rname, yvals, totval) in enumerate(rows_5y):
            r = r_idx + 1
            is_prof = "순영업익" in rname
            bg_c = self.c_tint_blue if is_prof else (self.c_white if r % 2 == 0 else self.c_box_bg)
            
            self._format_cell(table_s12.cell(r, 0), rname, font_size=9, bold=True, color=self.c_mck_navy if is_prof else self.c_charcoal, bg_color=bg_c)
            for idx, yv in enumerate(yvals):
                self._format_cell(table_s12.cell(r, idx+1), yv, font_size=9, bold=is_prof, color=self.c_red if is_prof else self.c_charcoal, bg_color=bg_c)
            self._format_cell(table_s12.cell(r, 6), totval, font_size=9.5, bold=True, color=self.c_mck_teal if is_prof else self.c_mck_navy, bg_color=bg_c)

        # 2. 좌측: 가맹점 출점 기대효과
        rect_l = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.72), Inches(5.9), Inches(4.35))
        rect_l.fill.solid()
        rect_l.fill.fore_color.rgb = self.c_white
        rect_l.line.color.rgb = self.c_line
        rect_l.line.width = Pt(1.2)
        
        top_bar_l = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.72), Inches(5.9), Inches(0.45))
        top_bar_l.fill.solid()
        top_bar_l.fill.fore_color.rgb = self.c_mck_navy
        top_bar_l.line.fill.background()
        
        tf_tbl = top_bar_l.text_frame
        p_tbl = tf_tbl.paragraphs[0]
        p_tbl.text = " 【 가맹점 출점 기대효과 및 핵심 경쟁력 】"
        p_tbl.font.size = Pt(11)
        p_tbl.font.bold = True
        p_tbl.font.color.rgb = self.c_white
        
        tb_l_body = s12.shapes.add_textbox(Inches(0.72), Inches(3.25), Inches(5.65), Inches(3.75))
        tf_lb = tb_l_body.text_frame
        tf_lb.word_wrap = True
        tf_lb.margin_left = tf_lb.margin_right = tf_lb.margin_top = tf_lb.margin_bottom = 0
        
        f_points = [
            ("1. 주간 유휴시간 제로 (정기 예약 중심 안정적 가동 체계)",
             f"• 일반 스크린골프 손님이 전무한 '평일 낮 10시~오후 5시' 유휴 시간대를 독점\n"
             f"• 반경 3km 내 {demo['senior_50_plus']/10000:.1f}만 50대 이상 시니어 및 여성 주부 동호회 4인 1팀 정기 리그 가동\n"
             f"• 비수기 및 날씨 영향을 받지 않는 사계절 정기 예약 중심 안정적 가동 안정성 확보"),
            (f"2. {site['rooms']}타석 플래그십 상위 20% 시장 선점",
             f"• 지역 내 소규모 매장 대비 {site['rooms']}타석 대규모 플래그십 시설 경쟁력 압도\n"
             f"• MYPARK 지역등급 추정 상위 20% 월매출 {comm['top_20_sales']//10000:,}만원 시장을 단독 선점 점유\n"
             f"• 카페형 휴게 라운지 및 파크골프 용품 샵 결합으로 객단가 및 체류시간 극대화"),
            ("3. 빠른 원금 회수 및 압도적 고수익성",
             f"• 직원 위탁 운영: 월 순영업익 약 {inv['owner_operated']['staff3_operating_profit']//10000:,}만원 (이익률 {inv['owner_operated']['staff3_operating_profit']/fin['monthly_scenarios']['moderate']['total_revenue']*100:.1f}%) / 회수 {inv['owner_operated']['staff3_payback_months']:.1f}개월\n"
             f"★ 창업주 직접 운영 시: 월 순영업익 {fin['monthly_scenarios']['moderate']['operating_profit']//10000:,}만원 / 단 {inv['payback_months_moderate']:.1f}개월({fmt_months(inv['payback_months_moderate'])}) 회수\n"
             f"• 손익분기점(BEP)은 타석당 하루 {inv['bep_turns_per_room']}회전(1일 {inv['bep_daily_users']}명) 수준으로 안정적인 가동 목표")
        ]
        for idx, (title, desc) in enumerate(f_points):
            p_t = tf_lb.add_paragraph() if idx > 0 else tf_lb.paragraphs[0]
            p_t.space_before = Pt(8) if idx > 0 else Pt(0)
            p_t.text = f"● {title}"
            p_t.font.size = Pt(10)
            p_t.font.bold = True
            p_t.font.color.rgb = self.c_mck_navy
            p_d = tf_lb.add_paragraph()
            p_d.space_before = Pt(3)
            p_d.text = desc
            p_d.font.size = Pt(8.6)
            p_d.font.color.rgb = self.c_charcoal

        if site.get('special_notes'):
            p_sn = tf_lb.add_paragraph()
            p_sn.space_before = Pt(6)
            p_sn.text = f"● 4. 고객 맞춤형 출점 전략 ('{site['special_notes']}')"
            p_sn.font.size = Pt(10)
            p_sn.font.bold = True
            p_sn.font.color.rgb = self.c_mck_navy
            p_snd = tf_lb.add_paragraph()
            p_snd.space_before = Pt(2)
            p_snd.text = f"• 고객 요청 특이사항('{site['special_notes']}')을 반영한 1:1 맞춤형 인테리어 및 운영 모델을 제안합니다."
            p_snd.font.size = Pt(8.6)
            p_snd.font.color.rgb = self.c_charcoal
            
        # 3. 우측: 건물주 및 상가 상생 활성화 효과
        rect_r = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.72), Inches(5.9), Inches(4.35))
        rect_r.fill.solid()
        rect_r.fill.fore_color.rgb = self.c_white
        rect_r.line.color.rgb = self.c_line
        rect_r.line.width = Pt(1.2)
        
        top_bar_r = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.72), Inches(5.9), Inches(0.45))
        top_bar_r.fill.solid()
        top_bar_r.fill.fore_color.rgb = self.c_mck_teal
        top_bar_r.line.fill.background()
        
        tf_tbr = top_bar_r.text_frame
        p_tbr = tf_tbr.paragraphs[0]
        p_tbr.text = " 【 건물주 및 상가 전체 상생 활성화 효과 】"
        p_tbr.font.size = Pt(11)
        p_tbr.font.bold = True
        p_tbr.font.color.rgb = self.c_white
        
        tb_r_body = s12.shapes.add_textbox(Inches(6.92), Inches(3.25), Inches(5.65), Inches(3.75))
        tf_rb = tb_r_body.text_frame
        tf_rb.word_wrap = True
        tf_rb.margin_left = tf_rb.margin_right = tf_rb.margin_top = tf_rb.margin_bottom = 0
        
        l_points = [
            ("1. 상가 내 구매력 높은 액티브 시니어 유동인구 앵커시설",
             f"• 보편 시나리오 기준 월 {m_scen['moderate']['monthly_users']:,}명 수준의 구매력 높은 50~70대 시니어 고객이 상가로 직접 유입\n"
             "• 4인 1팀 단체 이용 특성상 1층 카페, 음식점, 병의원 등 연계 소비 유발\n"
             "• 평일 낮 10~17시 상가 전체의 주간 공실 및 유휴 분위기를 완전히 반전"),
            ("2. 장기 안정적 우량 임차인 락인 (공실 리스크 영구 해소)",
             f"• 초기 설비 투자금 {fmt_eok(inv['total_capex'])}이 투입되는 실물 시설형 매장으로 5년 이상 장기 계약 유지\n"
             f"• 임대 운영 시 월 {site['monthly_rent']//10000:,}만원 수준의 체납 없는 안정적 수취 구조 (자가 운영 시 해당 없음)\n"
             f"• 상가 공실률 해소 및 앵커 테넌트 유치에 따른 건물 전체의 자산 가치(Cap Rate) 동반 상승"),
            ("3. 쾌적한 무소음·무진동·비음주 청정 친환경 체육시설",
             "• 주간 친목형 청정 체육시설로 음주·소음·흡연 부담이 없는 운영 형태\n"
             "• 상하층 입점 학원, 병원, 사무실과의 민원 마찰 전혀 없는 클린 테넌트\n"
             "• 쾌적한 건물 이미지 구축 및 시니어 친화 랜드마크 건물로 브랜딩 효과 극대화")
        ]
        for idx, (title, desc) in enumerate(l_points):
            p_t = tf_rb.add_paragraph() if idx > 0 else tf_rb.paragraphs[0]
            p_t.space_before = Pt(8) if idx > 0 else Pt(0)
            p_t.text = f"● {title}"
            p_t.font.size = Pt(10)
            p_t.font.bold = True
            p_t.font.color.rgb = self.c_mck_navy
            p_d = tf_rb.add_paragraph()
            p_d.space_before = Pt(3)
            p_d.text = desc
            p_d.font.size = Pt(8.6)
            p_d.font.color.rgb = self.c_charcoal

        self._add_source_footer(s12, "MYPARK 5-Year Financial Forecast & Final Strategic Recommendation")

        self.prs.save(self.filename)
        print(f"[PPTX GENERATED 12S] {self.filename}")
        return self.filename
