import os, sys
sys.path.insert(0, os.path.abspath('.'))
import os
import pptx
import docx
from tools.mypark_analyzer import MyParkReportGenerator

def test_suite():
    print('='*50)
    print('Starting MyPark Analyzer Test Suite')
    print('='*50)
    
    cases = [
        {
            'name': 'Case 1: Ilsan Pungdong',
            'address': '경기도 고양시 일산동구 숲속마을로 22',
            'building': 'Ilsan_Pungdong_Store',
            'rooms': 12,
            'rent': 5000000,
            'area': 120
        },
        {
            'name': 'Case 2: Ilsan Janghangdong',
            'address': '경기도 고양시 일산동구 장항동 736-6',
            'building': 'Ilsan_Janghang_Store',
            'rooms': 10,
            'rent': 4500000,
            'area': 100
        },
        {
            'name': 'Case 3: Songdo International City',
            'address': '인천광역시 연수구 하모니로177번길 49',
            'building': 'Songdo_Store',
            'rooms': 12,
            'rent': 5000000,
            'area': 150
        },
        {
            'name': 'Case 4: Yongin Suji Pungdeokcheon',
            'address': '경기도 용인시 수지구 풍덕천동 1082',
            'building': 'Yongin_Suji_Store',
            'rooms': 8,
            'rent': 4000000,
            'area': 85
        }
    ]
    
    gen = MyParkReportGenerator(output_dir='output')
    for c in cases:
        cname = c['name']
        print(f'\nRunning: {cname}')
        res = gen.analyze_and_generate(
            address=c['address'],
            building_name=c['building'],
            rooms=c['rooms'],
            monthly_rent=c['rent'],
            area_pyeong=c['area']
        )
        
        pptx_file = res['pptx_path']
        docx_file = res['docx_path']
        
        assert os.path.exists(pptx_file), f'Missing PPTX: {pptx_file}'
        assert os.path.exists(docx_file), f'Missing DOCX: {docx_file}'
        
        prs = pptx.Presentation(pptx_file)
        assert len(prs.slides) == 12, f'Expected 12 slides, got {len(prs.slides)}'
        
        doc = docx.Document(docx_file)
        assert len(doc.paragraphs) > 5, 'Too few paragraphs'
        assert len(doc.tables) >= 1, 'Too few tables'
        
        grade = res['grade']
        score = res['total_score']
        rev = res['total_revenue_moderate']
        op = res['operating_profit_moderate']
        pb = res['payback_months']
        pptx_base = os.path.basename(pptx_file)
        docx_base = os.path.basename(docx_file)
        
        print(f' PASS: {cname}')
        print(f'  - Grade: {grade} ({score} pts)')
        print(f'  - Monthly Revenue: {rev:,} KRW')
        print(f'  - Monthly Profit: {op:,} KRW')
        print(f'  - Payback: {pb} months')
        print(f'  - PPTX: {pptx_base} (12 slides)')
        print(f'  - DOCX: {docx_base} ({len(doc.paragraphs)} paras, {len(doc.tables)} tables)')
        
    print('\n' + '='*50)
    print('ALL 4 TEST CASES COMPLETED AND PASSED SUCCESSFULLY!')
    print('='*50)

if __name__ == '__main__':
    test_suite()
